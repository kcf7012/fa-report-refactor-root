"""LLM Evaluator — 整合 pptx 解析與 LLM 評估"""

from __future__ import annotations

import json
import re
from pathlib import Path
from typing import List, Optional

from pptx import Presentation

from ..domain.evaluation import Dimension, DimensionScore, EvaluationResult
from ..domain.suggestion import Improvement, Priority
from .base import LLMClient, LLMError, LLMResponse
from .prompts import SYSTEM_PROMPT, build_user_prompt


class LLMEvaluator:
    """使用 LLM 評估 FA 報告

    流程:
    1. 從 pptx 提取所有文字內容
    2. 呼叫 LLM 評估
    3. 將 LLM 回應(JSON)轉為 EvaluationResult
    """

    def __init__(self, llm_client: LLMClient):
        self.llm = llm_client

    def evaluate_pptx(self, pptx_path: Path) -> EvaluationResult:
        """評估 pptx 檔案"""
        content = self._extract_content(pptx_path)
        response = self.llm.complete(
            system_prompt=SYSTEM_PROMPT,
            user_prompt=build_user_prompt(content),
            json_mode=True,
            temperature=0.0,
        )
        return self._parse_response(response, source_file=str(pptx_path))

    def _extract_content(self, pptx_path: Path) -> str:
        """從 pptx 提取所有文字內容"""
        prs = Presentation(pptx_path)
        sections = []

        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_text.append(text)
                if hasattr(shape, "has_table") and shape.has_table:
                    table_text = self._extract_table(shape.table)
                    if table_text:
                        slide_text.append(table_text)

            if slide_text:
                sections.append(f"=== Slide {i} ===\n" + "\n".join(slide_text))

        return "\n\n".join(sections)

    def _extract_table(self, table) -> str:
        """提取表格內容為文字"""
        rows = []
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            rows.append(row_text)
        return "\n".join(rows)

    def _parse_response(self, response: LLMResponse, source_file: str = "") -> EvaluationResult:
        """解析 LLM 回應為 EvaluationResult"""
        content = response.content.strip()

        # 嘗試從可能的 markdown code block 中提取 JSON
        json_match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", content, re.DOTALL)
        if json_match:
            content = json_match.group(1)

        try:
            data = json.loads(content)
        except json.JSONDecodeError as e:
            raise LLMError(
                f"LLM 回應不是有效 JSON:{e}\n"
                f"內容:{content[:500]}"
            ) from e

        return self._from_dict(data, response, source_file)

    def _from_dict(
        self,
        data: dict,
        response: LLMResponse,
        source_file: str,
    ) -> EvaluationResult:
        """從 LLM 回應的 dict 建立 EvaluationResult"""
        # 解析維度分數
        dimensions = []
        dims_data = data.get("dimension_scores", data.get("dimensions", {}))

        for name, dim_data in dims_data.items():
            if isinstance(dim_data, dict):
                dim = self._match_dimension(name)
                dimensions.append(
                    DimensionScore(
                        name=dim,
                        score=float(dim_data.get("score", 0)),
                        weight=int(dim_data.get("weight", dim.weight)),
                        comment=dim_data.get("comment", ""),
                    )
                )
            elif isinstance(dim_data, (int, float)):
                dim = self._match_dimension(name)
                dimensions.append(
                    DimensionScore(
                        name=dim,
                        score=float(dim_data),
                        weight=dim.weight,
                    )
                )

        # 解析改進建議
        improvements = []
        for imp in data.get("improvements", []):
            improvements.append(self._parse_improvement(imp))

        return EvaluationResult(
            total_score=float(data.get("total_score", 0)),
            grade=data.get("grade", "F"),
            dimensions=dimensions,
            summary=data.get("summary", ""),
            strengths=data.get("strengths", []),
            source_file=source_file,
            token_usage={
                "prompt_tokens": response.prompt_tokens,
                "completion_tokens": response.completion_tokens,
                "total_tokens": response.total_tokens,
            },
        )

    def _parse_improvement(self, imp) -> Improvement:
        """解析改進建議項目"""
        if isinstance(imp, str):
            return Improvement.from_text(imp)
        return Improvement.from_dict(imp)

    def _match_dimension(self, name: str) -> Dimension:
        """從名稱匹配 Dimension"""
        for dim in Dimension:
            if dim.value == name or name in dim.value:
                return dim
        # fallback
        raise ValueError(f"無法識別維度: {name}")