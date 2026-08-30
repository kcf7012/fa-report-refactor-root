"""母片保護機制 — 確保改善過程不破壞母片"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import List

from pptx import Presentation


@dataclass
class MasterSnapshot:
    """母片的不可變快照"""

    masters_xml: List[str] = field(default_factory=list)
    layouts_xml: List[str] = field(default_factory=list)
    layout_names: List[str] = field(default_factory=list)
    image_count: int = 0

    def diff(self, other: "MasterSnapshot") -> List[str]:
        """比較兩個快照,回傳差異清單"""
        violations = []
        if self.masters_xml != other.masters_xml:
            violations.append("母片 XML 被修改")
        if self.layouts_xml != other.layouts_xml:
            violations.append(f"Layout XML 被修改({len(self.layouts_xml)} → {len(other.layouts_xml)})")
        if self.layout_names != other.layout_names:
            violations.append(
                f"Layout 數量/名稱變更:{len(self.layout_names)} → {len(other.layout_names)}"
            )
        if self.image_count != other.image_count:
            violations.append(
                f"圖片數量變更:{self.image_count} → {other.image_count}"
            )
        return violations


class MasterProtector:
    """三層保護機制

    1. 改善開始時擷取母片快照
    2. 每次改善動作前檢查必要條件
    3. 改善結束後驗證母片未被修改
    """

    def __init__(self, prs: Presentation):
        self.snapshot = self._capture(prs)

    def _capture(self, prs: Presentation) -> MasterSnapshot:
        """擷取所有不可變狀態"""
        return MasterSnapshot(
            masters_xml=[m.element.xml for m in prs.slide_masters],
            layouts_xml=[l.element.xml for l in prs.slide_layouts],
            layout_names=[l.name for l in prs.slide_layouts],
            image_count=self._count_master_images(prs),
        )

    def _count_master_images(self, prs: Presentation) -> int:
        """計算母片相關圖片數量"""
        count = 0
        for master in prs.slide_masters:
            for shape in master.shapes:
                if shape.shape_type == 13:  # PICTURE
                    count += 1
        # 也算上 layouts 的圖片
        for layout in prs.slide_layouts:
            for shape in layout.shapes:
                if shape.shape_type == 13:
                    count += 1
        return count

    def verify_unchanged(self, prs: Presentation) -> None:
        """驗證母片從改善開始到結束都未被修改

        Raises:
            MasterProtectionError: 母片被修改時
        """
        current = self._capture(prs)
        violations = self.snapshot.diff(current)
        if violations:
            raise MasterProtectionError(
                "母片保護失敗!\以下項目被修改:\n" + "\n".join(f"  - {v}" for v in violations)
            )

    def assert_can_add_slide(self, prs: Presentation, layout_name: str) -> None:
        """確認可以使用指定的 layout 新增投影片

        Raises:
            MasterProtectionError: layout 不存在
        """
        existing_names = [l.name for l in prs.slide_layouts]
        if layout_name not in existing_names:
            raise MasterProtectionError(
                f"Layout '{layout_name}' 不存在!\n"
                f"現有 layouts:{existing_names}\n"
                f"禁止建立新 layout,必須使用既有 layout。"
            )


class MasterProtectionError(Exception):
    """母片保護失敗時拋出"""

    pass