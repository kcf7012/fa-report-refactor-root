"""Summary 強化 — 展開為多張投影片"""

from __future__ import annotations

from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from ..domain.evaluation import EvaluationResult
from ..domain.suggestion import Improvement
from ..layout.selector import find_content_layout
from ..visuals import (
    ELAN_BLUE,
    ELAN_GREEN,
    ELAN_ORANGE,
    ELAN_RED,
    ProgressBarGenerator,
)


def enhance_summary_section(
    prs: Presentation,
    evaluation: EvaluationResult,
    improvements: List[Improvement],
) -> None:
    """強化 Summary 區塊

    策略:
    - 保留原 Summary 投影片不動
    - 注入 Executive Summary 與 Key Improvements
    - 若有空間,加入 6 維度評分視覺化
    """
    summary_idx = _find_summary_index(prs)
    if summary_idx == -1:
        summary_idx = len(prs.slides) - 1
    if summary_idx < 0:
        return

    slide = prs.slides[summary_idx]

    # 注入 6 維度評分進度條(若有資料)
    if evaluation.dimensions:
        _add_dimension_progress(slide, evaluation)

    # 注入 Executive Summary
    _add_executive_summary(slide, evaluation)

    # 注入 Key Improvements
    _add_key_improvements(slide, improvements, evaluation)

    # 注入分析優點
    _add_strengths(slide, evaluation)


def _find_summary_index(prs: Presentation) -> int:
    """尋找 Summary/總結 投影片"""
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text
                if any(kw in text for kw in ["Summary", "總結", "Executive Summary"]):
                    return i
    return -1


def _add_executive_summary(slide, evaluation: EvaluationResult) -> None:
    """加入 Executive Summary 文字框"""
    textbox = slide.shapes.add_textbox(
        Inches(7.5), Inches(3.0), Inches(4.5), Inches(1.5)
    )
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = evaluation.summary or "報告分析詳實,建議補充統計數據以強化結論。"
    p.font.size = Pt(11)


def _add_key_improvements(
    slide, improvements: List[Improvement], evaluation: EvaluationResult
) -> None:
    """加入 Key Improvements Required 文字框"""
    textbox = slide.shapes.add_textbox(
        Inches(7.5), Inches(4.6), Inches(4.5), Inches(2.0)
    )
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Improvements Required"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 0, 0)

    if not improvements:
        improvements = [
            Improvement(
                priority=improvements[0].priority if improvements else __import__(
                    "fa_improver.domain.suggestion", fromlist=["Priority"]
                ).Priority.MEDIUM,
                item="改善建議",
                suggestion="依評核建議補強缺失項目",
            )
        ]

    for imp in improvements[:3]:
        p = tf.add_paragraph()
        p.text = f"• {imp.suggestion}"
        p.font.size = Pt(10)


def _add_strengths(slide, evaluation: EvaluationResult) -> None:
    """加入分析優點與成功驗證"""
    if not evaluation.strengths:
        return

    textbox = slide.shapes.add_textbox(
        Inches(1.6), Inches(4.6), Inches(3.6), Inches(2.0)
    )
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "分析優點與成功驗證"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 112, 192)

    for s in evaluation.strengths[:5]:
        p = tf.add_paragraph()
        p.text = f"✓ {s}"
        p.font.size = Pt(10)


def _add_dimension_progress(slide, evaluation: EvaluationResult) -> None:
    """加入 6 維度評分進度條"""
    if not evaluation.dimensions:
        return

    # 根據分數決定顏色
    def color_for_score(score: float):
        if score >= 85:
            return ELAN_GREEN
        if score >= 70:
            return ELAN_BLUE
        if score >= 50:
            return ELAN_ORANGE
        return ELAN_RED

    items = [
        {
            "label": dim_score.name.value,
            "value": dim_score.score,
            "max_value": 100,
            "color": color_for_score(dim_score.score),
        }
        for dim_score in evaluation.dimensions
    ]

    # 放在原有內容下方(不重疊)
    gen = ProgressBarGenerator(
        slide,
        left=1.0,
        top=5.2,
        width=8.0,
        height=2.0,
    )
    gen.generate(items)