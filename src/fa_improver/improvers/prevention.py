"""新增改善對策投影片"""

from __future__ import annotations

from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from ..domain.evaluation import EvaluationResult
from ..domain.suggestion import Improvement
from ..layout.selector import find_content_layout


def add_prevention_measures_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    improvements: List[Improvement],
) -> None:
    """新增長期預防措施與改善對策投影片"""
    layout = find_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    title = _get_or_create_title(slide)
    title.text_frame.text = "長期預防措施與改善對策"

    body = _get_or_create_body(slide)
    tf = body.text_frame
    tf.clear()

    # 擬議改善對策項目
    p = tf.paragraphs[0]
    p.text = "擬議改善對策項目:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 112, 192)

    # 從 improvements 抽取建議
    if improvements:
        for imp in improvements[:3]:
            p = tf.add_paragraph()
            p.text = imp.suggestion
            p.font.size = Pt(14)

    # 標準化與監測計畫
    p = tf.add_paragraph()
    p.text = "\n[標準化與監測計畫]"
    p.font.bold = True

    for item in [
        "建立入料檢驗 (IQC) SOP 與測試閾值",
        "導入自動化監測設備於生產線",
        "將此案例納入知識管理資料庫以利後續追蹤",
    ]:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)


def _get_or_create_title(slide):
    if slide.shapes.title:
        return slide.shapes.title
    for shape in slide.shapes:
        if "title" in shape.name.lower():
            return shape
    return slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))


def _get_or_create_body(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            return shape
    return slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))