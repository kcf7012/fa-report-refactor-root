"""數據與證據支持 Improver

針對「數據與證據支持」維度(權重 20%)。
改善內容:
- 對照組 vs 異常品 數據對照表
- 圖片品質檢查清單
- 數據追溯性指引
"""

from __future__ import annotations

from typing import Optional

from pptx import Presentation
from pptx.util import Inches

from ..domain.evaluation import DimensionScore, EvaluationResult
from ..layout.selector import find_content_layout
from ..visuals import (
    ChecklistGenerator,
    ComparisonTableGenerator,
    ELAN_BLUE,
    ELAN_GREEN,
    ELAN_ORANGE,
    ELAN_RED,
)


def add_evidence_checklist_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    dimension: Optional[DimensionScore] = None,
) -> None:
    """新增數據與證據支持投影片

    Args:
        prs: 簡報物件
        evaluation: 評估結果
        dimension: 該維度的評分(可選)
    """
    layout = find_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    # 標題
    title = _get_or_create_title(slide)
    title.text_frame.text = "數據與證據支持"

    # 1. 對照組 vs 異常品 數據對照表
    _add_comparison_data_table(slide)

    # 2. 圖片品質與數據追溯性檢查清單
    _add_evidence_checklist(slide)


def _add_comparison_data_table(slide) -> None:
    """加入對照組 vs 異常品數據對照表"""
    gen = ComparisonTableGenerator(
        slide,
        left=0.5,
        top=1.4,
        width=8.5,
        height=2.2,
    )
    gen.generate(
        {
            "headers": ["測試項目", "DVT 正常品", "PVT 異常品", "Spec 範圍", "判定"],
            "rows": [
                ["ESD HBM (±2kV)", "通過", "?", "±2kV", "需補測"],
                ["I/O 對地阻抗", "16MΩ", "5.7KΩ", ">1MΩ", "FAIL"],
                ["VH/VOUT 電壓", "正常", "0V", "0.4-0.6V", "FAIL"],
                ["二極體特性", "0.43V", "0V", "0.4-0.6V", "FAIL"],
                ["FW 讀取", "正確", "正確", "100%", "PASS"],
                ["外觀檢查", "正常", "正常", "無損傷", "PASS"],
            ],
        }
    )


def _add_evidence_checklist(slide) -> None:
    """加入圖片品質與數據追溯性檢查清單"""
    gen = ChecklistGenerator(
        slide,
        left=0.5,
        top=3.9,
        width=8.5,
        height=2.8,
    )
    gen.generate(
        [
            {
                "text": "圖片解析度 ≥ 1000x(建議 ≥ 2000x 以利 IC Marking 辨識)",
                "checked": False,
                "color": ELAN_RED,
            },
            {
                "text": "對焦清晰,IC 標籤、Marking 可清楚辨識",
                "checked": False,
                "color": ELAN_RED,
            },
            {
                "text": "含比例尺(scale bar),方便評估損傷範圍",
                "checked": False,
                "color": ELAN_ORANGE,
            },
            {
                "text": "異常點明確標註(箭頭、方框、文字說明)",
                "checked": False,
                "color": ELAN_ORANGE,
            },
            {
                "text": "對照組(Golden Sample)與異常品並列比較",
                "checked": False,
                "color": ELAN_BLUE,
            },
            {
                "text": "每個量化數據有來源追溯(儀器型號、測試條件、日期)",
                "checked": False,
                "color": ELAN_BLUE,
            },
            {
                "text": "統計顯著性(p-value、CI)支援結論(非僅描述)",
                "checked": False,
                "color": ELAN_GREEN,
            },
        ]
    )


def _get_or_create_title(slide):
    if slide.shapes.title:
        return slide.shapes.title
    for shape in slide.shapes:
        if "title" in shape.name.lower():
            return shape
    return slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))