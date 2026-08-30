"""新增 FA 基本資訊投影片"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.util import Inches

from ..domain.evaluation import EvaluationResult
from ..layout.selector import find_content_layout
from ..parsers.filename_parser import FilenameInfo

if TYPE_CHECKING:
    pass


def add_basic_info_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    filename_info: FilenameInfo,
) -> None:
    """新增 FA 基本資訊投影片"""
    layout = find_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    # 標題
    title_shape = _get_or_create_title(slide)
    title_shape.text_frame.text = "FA 基本資訊"

    # 內容
    info_items = [
        ("FA 編號", filename_info.to_fa_id()),
        ("負責工程師", "ELAN FAE"),
        ("客戶", filename_info.customer or "N/A"),
        ("專案名稱", filename_info.project or "N/A"),
        ("報告日期", filename_info.date or "N/A"),
        ("失效數量", "依評核建議補充填寫"),
        ("批號 (Lot No.)", "依評核建議補充填寫"),
    ]

    body = _get_or_create_body(slide)
    tf = body.text_frame
    tf.clear()
    for label, value in info_items:
        p = tf.add_paragraph()
        p.text = f"{label}: {value}"
        p.font.size = _PT(14)

    # 優化建議項目(從 comment 抽取)
    if evaluation.dimensions:
        dim = next((d for d in evaluation.dimensions if d.name.value == "基本資訊完整性"), None)
        if dim and dim.comment:
            p = tf.add_paragraph()
            p.text = "\n[優化建議項目]"
            p.font.bold = True
            p.font.color.rgb = _COLOR(255, 0, 0)

            sub = tf.add_paragraph()
            sub.text = f"• {dim.comment}"
            sub.font.size = _PT(12)


def _get_or_create_title(slide):
    if slide.shapes.title:
        return slide.shapes.title
    for shape in slide.shapes:
        if "title" in shape.name.lower():
            return shape
    return slide.shapes.add_textbox(_Inches(0.5), _Inches(0.3), _Inches(9), _Inches(1))


def _get_or_create_body(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            return shape
    return slide.shapes.add_textbox(_Inches(0.5), _Inches(1.5), _Inches(9), _Inches(5))


# Helper aliases(縮短程式碼)
_Inches = Inches


def _PT(size: int):
    from pptx.util import Pt

    return Pt(size)


def _COLOR(r: int, g: int, b: int):
    from pptx.dml.color import RGBColor

    return RGBColor(r, g, b)