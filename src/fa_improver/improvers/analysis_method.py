"""分析方法與流程 Improver

針對「分析方法與流程」維度(權重 20%)。
改善內容:
- 8D 流程檢查清單(D1-D8)
- 分析方法對照表(SEM/FIB/X-ray 適用場景)
- 實驗設計 SOP 範本
"""

from __future__ import annotations

from typing import Optional

from pptx import Presentation
from pptx.util import Inches

from ..domain.evaluation import DimensionScore, EvaluationResult
from ..layout.selector import find_content_layout
from ..visuals import (
    ChecklistGenerator,
    ComparisonTableGenerator,
    ELAN_BLUE,
    ELAN_GREEN,
    ELAN_GRAY,
    ELAN_ORANGE,
    ELAN_RED,
)


def add_analysis_method_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    dimension: Optional[DimensionScore] = None,
) -> None:
    """新增分析方法與流程投影片

    Args:
        prs: 簡報物件
        evaluation: 評估結果
        dimension: 該維度的評分(可選)
    """
    layout = find_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    # 標題
    title = _get_or_create_title(slide)
    title.text_frame.text = "分析方法與流程"

    # 1. 8D 流程檢查清單
    _add_8d_checklist(slide)

    # 2. 分析方法對照表
    _add_method_comparison_table(slide)


def _add_8d_checklist(slide) -> None:
    """加入 8D 流程檢查清單"""
    gen = ChecklistGenerator(
        slide,
        left=0.5,
        top=1.4,
        width=4.0,
        height=4.5,
    )
    gen.generate(
        [
            {"text": "D1: 建立團隊(FAE/QRA/PM)", "checked": False, "color": ELAN_BLUE},
            {"text": "D2: 描述問題(5W2H)", "checked": False, "color": ELAN_BLUE},
            {"text": "D3: 圍堵行動(Containment)", "checked": False, "color": ELAN_ORANGE},
            {"text": "D4: 根因分析(5-Why/Fishbone)", "checked": False, "color": ELAN_RED},
            {"text": "D5: 永久對策(PCA)", "checked": False, "color": ELAN_RED},
            {"text": "D6: 實施與驗證", "checked": False, "color": ELAN_ORANGE},
            {"text": "D7: 再發防止(SOP/防呆)", "checked": False, "color": ELAN_BLUE},
            {"text": "D8: 團隊表揚與經驗傳承", "checked": False, "color": ELAN_GREEN},
        ]
    )


def _add_method_comparison_table(slide) -> None:
    """加入分析方法對照表"""
    gen = ComparisonTableGenerator(
        slide,
        left=5.0,
        top=1.4,
        width=4.5,
        height=4.5,
    )
    gen.generate(
        {
            "headers": ["方法", "適用場景", "備註"],
            "rows": [
                ["光學檢查", "外部損傷、燒毀痕跡", "快速、低成本"],
                ["SEM", "晶片內部結構", "需真空環境"],
                ["FIB", "橫截面分析", "破壞性測試"],
                ["X-ray", "封裝裂縫、空洞", "非破壞性"],
                ["Decap", "開蓋檢查", "化學蝕刻"],
                ["TDR/示波器", "訊號完整性", "時域反射"],
                ["XRD/EDX", "材料分析", "成分鑑定"],
                ["ESD 測試", "靜電耐受", "HBM/CDM 模型"],
            ],
        }
    )


def _get_or_create_title(slide):
    if slide.shapes.title:
        return slide.shapes.title
    for shape in slide.shapes:
        if "title" in shape.name.lower():
            return shape
    return slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))