"""新增根因分析相關投影片"""

from __future__ import annotations

from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from ..domain.evaluation import EvaluationResult
from ..layout.selector import find_content_layout


def add_statistical_analysis_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    suggestions: List[str],
    variant: str = "statistical",
) -> None:
    """新增根因分析投影片

    variant:
    - "5_why": 5-Why 推導流程
    - "statistical": 統計驗證方法
    """
    layout = find_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    title = _get_or_create_title(slide)
    title.text_frame.text = (
        "5-Why 根因推導" if variant == "5_why" else "根因驗證及統計分析"
    )

    # 內容
    if not suggestions:
        suggestions = ["建議加強對照組設定與數據統計驗證以支撐根因發現。"]

    body = _get_or_create_body(slide)
    tf = body.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "針對問題點之深度分析建議:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 112, 192)

    for sug in suggestions[:4]:
        p = tf.add_paragraph()
        p.text = sug
        p.font.size = Pt(14)

    # 建議執行動作
    p = tf.add_paragraph()
    p.text = "\n[建議執行動作]"
    p.font.bold = True

    for action in [
        "設定 DVT 正常品 vs PVT 異常品之對照組",
        "使用獨立樣本 t 檢定驗證參數顯著性 (p < 0.05)",
        "確保統計證據支持最終提到的根本原因",
    ]:
        p = tf.add_paragraph()
        p.text = f"• {action}"
        p.font.size = Pt(12)


def _get_or_create_title(slide):
    if slide.shapes.title:
        return slide.shapes.title
    for shape in slide.shapes:
        if "title" in shape.name.lower():
            return shape
    return slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))


def _get_or_create_body(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            return shape
    return slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))