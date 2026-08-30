"""問題描述與定義 Improver

針對「問題描述與定義」維度(權重 15%)。
改善內容:
- 失效現象 vs 失效模式 對照表
- 問題範圍量化(失效率、影響數量)
- 客戶影響評估
"""

from __future__ import annotations

from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from ..domain.evaluation import DimensionScore, EvaluationResult
from ..layout.selector import find_content_layout
from ..visuals import (
    ChecklistGenerator,
    ComparisonTableGenerator,
    ELAN_BLUE,
    ELAN_GREEN,
    ELAN_ORANGE,
    ELAN_RED,
)


def add_problem_definition_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    dimension: Optional[DimensionScore] = None,
) -> None:
    """新增問題描述與定義投影片

    Args:
        prs: 簡報物件
        evaluation: 評估結果
        dimension: 該維度的評分(可選,用於針對性改善)
    """
    layout = find_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    # 標題
    title = _get_or_create_title(slide)
    title.text_frame.text = "問題描述與失效定義"

    # 1. 失效現象 vs 失效模式對照表
    _add_phenomenon_vs_mode_table(slide)

    # 2. 問題範圍量化檢查清單
    _add_quantification_checklist(slide)


def _add_phenomenon_vs_mode_table(slide) -> None:
    """加入失效現象 vs 失效模式對照表"""
    gen = ComparisonTableGenerator(
        slide,
        left=0.5,
        top=1.4,
        width=8.5,
        height=2.0,
    )
    gen.generate(
        {
            "headers": ["項目", "目前報告常見缺失", "建議補充內容"],
            "rows": [
                [
                    "失效現象 (Phenomenon)",
                    "僅描述「通訊失敗」等表面現象",
                    "完整描述:電壓/電流/波形/時序/外觀",
                ],
                [
                    "失效模式 (Failure Mode)",
                    "未明確定義失效機制",
                    "對應失效機制:開路/短路/漏電/漂移/功能失效",
                ],
                [
                    "失效位置",
                    "只說「IC 損壞」",
                    "精確位置:腳位/區塊/層次(晶圓/封裝/PCB)",
                ],
            ],
        }
    )


def _add_quantification_checklist(slide) -> None:
    """加入問題範圍量化檢查清單"""
    gen = ChecklistGenerator(
        slide,
        left=0.5,
        top=3.7,
        width=8.5,
        height=3.0,
    )
    gen.generate(
        [
            {
                "text": "失效率(PPM 或 %):本次失效佔總出貨量的比例",
                "checked": False,
                "color": ELAN_RED,
            },
            {
                "text": "影響產品數量:目前庫存/在製品/已出貨的受影響數量",
                "checked": False,
                "color": ELAN_ORANGE,
            },
            {
                "text": "影響時間範圍:失效首次發生日期 → 目前(持續中/已結案)",
                "checked": False,
                "color": ELAN_ORANGE,
            },
            {
                "text": "客戶影響評估:客戶端失效比例、生產線停線時間、退貨/召回成本",
                "checked": False,
                "color": ELAN_RED,
            },
            {
                "text": "失效嚴重性分級:Critical(安全)/Major(功能)/Minor(性能降級)",
                "checked": False,
                "color": ELAN_RED,
            },
            {
                "text": "批量 vs 個案:是否為批次性失效(同批多顆)或個案(單顆)",
                "checked": False,
                "color": ELAN_BLUE,
            },
        ]
    )


def _get_or_create_title(slide):
    if slide.shapes.title:
        return slide.shapes.title
    for shape in slide.shapes:
        if "title" in shape.name.lower():
            return shape
    return slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))