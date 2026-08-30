"""視覺元素生成器的抽象基類與實作"""

from __future__ import annotations

from abc import ABC, abstractmethod
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .colors import (
    ELAN_BLUE,
    ELAN_GREEN,
    ELAN_GRAY,
    ELAN_LIGHT_BLUE,
    ELAN_LIGHT_GRAY,
    ELAN_ORANGE,
    ELAN_RED,
)


class VisualGenerator(ABC):
    """視覺元素生成器抽象基類

    所有生成器都接受:
    - slide: 目標投影片
    - left/top/width/height: 位置與大小
    - content: 內容資料(由各生成器自行定義)
    """

    def __init__(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
    ):
        self.slide = slide
        self.left = Inches(left)
        self.top = Inches(top)
        self.width = Inches(width)
        self.height = Inches(height)

    @abstractmethod
    def generate(self, content) -> None:
        """生成視覺元素"""

    def _set_text(
        self,
        shape,
        text: str,
        *,
        size: int = 12,
        bold: bool = False,
        color: RGBColor = None,
        align=None,
    ):
        """設定 shape 文字內容與樣式"""
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)

        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


class ChecklistGenerator(VisualGenerator):
    """Checkbox 列表生成器

    內容格式:List[str] 或 List[dict(text, checked, priority)]
    """

    def generate(self, content: List) -> None:
        from pptx.oxml.ns import qn

        items = self._normalize(content)

        # 計算每個項目的高度
        item_count = len(items)
        if item_count == 0:
            return
        item_height = self.height / item_count

        for i, item in enumerate(items):
            top = self.top + item_height * i
            # checkbox 方塊
            box = self.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                self.left,
                top,
                Inches(0.25),
                Inches(0.25),
            )
            box.line.color.rgb = ELAN_GRAY
            box.fill.solid()
            box.fill.fore_color.rgb = ELAN_LIGHT_GRAY if not item["checked"] else ELAN_GREEN
            self._set_text(box, "✓" if item["checked"] else "☐", size=14, bold=True, align=PP_ALIGN.CENTER)

            # 文字
            text_box = self.slide.shapes.add_textbox(
                self.left + Inches(0.35),
                top,
                self.width - Inches(0.35),
                item_height,
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = item["text"]
            run.font.size = Pt(11)
            run.font.color.rgb = item["color"]

    def _normalize(self, content: List) -> List[dict]:
        """標準化輸入格式"""
        result = []
        for item in content:
            if isinstance(item, str):
                result.append({"text": item, "checked": False, "color": ELAN_BLUE})
            elif isinstance(item, dict):
                result.append(
                    {
                        "text": item.get("text", ""),
                        "checked": item.get("checked", False),
                        "color": item.get("color", ELAN_BLUE),
                    }
                )
        return result


class FlowDiagramGenerator(VisualGenerator):
    """流程圖生成器

    內容格式:List[str] (步驟名稱)
    或 List[dict(name, status)]
    """

    def generate(self, content: List) -> None:
        steps = self._normalize(content)
        if not steps:
            return

        step_count = len(steps)
        # 每個步驟的水平寬度
        step_width = self.width / step_count

        for i, step in enumerate(steps):
            left = self.left + step_width * i
            # 步驟方塊
            box = self.slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left + Inches(0.1),
                self.top,
                step_width - Inches(0.2),
                self.height,
            )
            box.fill.solid()
            box.fill.fore_color.rgb = step["color"]
            box.line.color.rgb = step["color"]

            self._set_text(
                box,
                step["name"],
                size=11,
                bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER,
            )

            # 箭頭(連到下一個)
            if i < step_count - 1:
                arrow = self.slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    left + step_width - Inches(0.15),
                    self.top + self.height / 2 - Inches(0.1),
                    Inches(0.2),
                    Inches(0.2),
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = ELAN_LIGHT_GRAY
                arrow.line.fill.background()

    def _normalize(self, content: List) -> List[dict]:
        """標準化輸入,根據 status 給顏色"""
        result = []
        for item in content:
            if isinstance(item, str):
                result.append(
                    {"name": item, "status": "done", "color": ELAN_BLUE}
                )
            elif isinstance(item, dict):
                status = item.get("status", "done")
                color = {
                    "done": ELAN_GREEN,
                    "current": ELAN_ORANGE,
                    "missing": ELAN_RED,
                    "pending": ELAN_LIGHT_BLUE,
                }.get(status, ELAN_BLUE)
                result.append(
                    {"name": item.get("name", ""), "status": status, "color": color}
                )
        return result


class ComparisonTableGenerator(VisualGenerator):
    """對照表生成器(原生 PowerPoint table)

    內容格式:List[List[str]] 或 dict(headers, rows)
    """

    def generate(self, content) -> None:
        if isinstance(content, dict):
            headers = content.get("headers", [])
            rows = content.get("rows", [])
        else:
            # 假設第一列是 headers
            headers = content[0] if content else []
            rows = content[1:] if len(content) > 1 else []

        if not headers:
            return

        # 建立 table
        table_shape = self.slide.shapes.add_table(
            len(rows) + 1,
            len(headers),
            self.left,
            self.top,
            self.width,
            self.height,
        )
        table = table_shape.table

        # 填入 headers
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = ELAN_BLUE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    run.font.bold = True
                    run.font.size = Pt(11)

        # 填入 rows
        for i, row in enumerate(rows, start=1):
            for j, value in enumerate(row):
                if j >= len(headers):
                    break
                cell = table.cell(i, j)
                cell.text = str(value)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)


class ProgressBarGenerator(VisualGenerator):
    """進度條生成器

    內容格式:List[dict(label, value, max_value, color)]
    """

    def generate(self, content: List[dict]) -> None:
        items = self._normalize(content)
        if not items:
            return

        bar_height = (self.height - Inches(0.3)) / len(items)
        label_width = Inches(1.5)
        bar_area_width = self.width - label_width - Inches(0.5)  # 留空間給數值文字

        for i, item in enumerate(items):
            top = self.top + bar_height * i

            # 標籤
            label_box = self.slide.shapes.add_textbox(
                self.left,
                top,
                label_width,
                bar_height,
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = item["label"]
            run.font.size = Pt(10)
            run.font.bold = True

            # 進度條背景(灰色)
            bg_height = bar_height - Inches(0.1)
            bg_box = self.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                self.left + label_width,
                top + Inches(0.05),
                bar_area_width,
                bg_height,
            )
            bg_box.fill.solid()
            bg_box.fill.fore_color.rgb = ELAN_LIGHT_GRAY
            bg_box.line.fill.background()

            # 進度條(實際值)
            ratio = max(0, min(1, item["value"] / item["max_value"]))
            fg_width = int(bar_area_width * ratio)
            if fg_width > 0:
                fg_box = self.slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    self.left + label_width,
                    top + Inches(0.05),
                    fg_width,
                    bg_height,
                )
                fg_box.fill.solid()
                fg_box.fill.fore_color.rgb = item["color"]
                fg_box.line.fill.background()

            # 數值文字
            value_box = self.slide.shapes.add_textbox(
                self.left + label_width + bar_area_width + Inches(0.05),
                top,
                Inches(0.45),
                bar_height,
            )
            tf = value_box.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"{item['value']:.0f}/{item['max_value']:.0f}"
            run.font.size = Pt(9)
            run.font.color.rgb = item["color"]
            run.font.bold = True

    def _normalize(self, content: List) -> List[dict]:
        """標準化輸入"""
        result = []
        for item in content:
            if isinstance(item, dict):
                value = float(item.get("value", 0))
                max_value = float(item.get("max_value", 100))
                color = item.get("color", ELAN_BLUE)
                result.append(
                    {
                        "label": item.get("label", ""),
                        "value": value,
                        "max_value": max_value,
                        "color": color,
                    }
                )
        return result


class TimelineGenerator(VisualGenerator):
    """時間軸生成器

    內容格式:List[dict(label, timeframe, color)]
    """

    def generate(self, content: List[dict]) -> None:
        phases = self._normalize(content)
        if not phases:
            return

        phase_count = len(phases)
        phase_width = self.width / phase_count

        # 時間軸主線
        axis = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.left,
            self.top + self.height / 2 - Inches(0.05),
            self.width,
            Inches(0.1),
        )
        axis.fill.solid()
        axis.fill.fore_color.rgb = ELAN_GRAY
        axis.line.fill.background()

        for i, phase in enumerate(phases):
            left = self.left + phase_width * i
            center_x = left + phase_width / 2

            # 階段節點(圓形)
            node = self.slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                center_x - Inches(0.2),
                self.top + self.height / 2 - Inches(0.2),
                Inches(0.4),
                Inches(0.4),
            )
            node.fill.solid()
            node.fill.fore_color.rgb = phase["color"]
            node.line.color.rgb = phase["color"]
            self._set_text(node, str(i + 1), size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

            # 標籤
            label_box = self.slide.shapes.add_textbox(
                left,
                self.top + self.height / 2 + Inches(0.25),
                phase_width,
                Inches(0.4),
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = phase["label"]
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = phase["color"]

            # 時程
            if phase.get("timeframe"):
                time_box = self.slide.shapes.add_textbox(
                    left,
                    self.top + self.height / 2 + Inches(0.55),
                    phase_width,
                    Inches(0.3),
                )
                tf = time_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = phase["timeframe"]
                run.font.size = Pt(9)
                run.font.color.rgb = ELAN_GRAY

    def _normalize(self, content: List) -> List[dict]:
        """標準化輸入"""
        result = []
        for i, item in enumerate(content):
            if isinstance(item, str):
                result.append(
                    {"label": item, "timeframe": "", "color": ELAN_BLUE}
                )
            elif isinstance(item, dict):
                color = item.get("color", ELAN_BLUE)
                result.append(
                    {
                        "label": item.get("label", ""),
                        "timeframe": item.get("timeframe", ""),
                        "color": color,
                    }
                )
        return result