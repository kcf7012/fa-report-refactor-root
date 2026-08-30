"""新 3 個 Improvers 測試(Phase 4.5)"""

import pytest
from pptx import Presentation

from fa_improver.improvers.problem_definition import add_problem_definition_slide
from fa_improver.improvers.analysis_method import add_analysis_method_slide
from fa_improver.improvers.evidence_checklist import add_evidence_checklist_slide


class TestProblemDefinitionImprover:
    """問題描述與定義 Improver 測試"""

    def test_add_slide_basic(self):
        """基本新增測試"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        from fa_improver.domain.evaluation import EvaluationResult

        evaluation = EvaluationResult(total_score=50.0, grade="F")
        add_problem_definition_slide(prs, evaluation)

        # 應新增一張投影片
        assert len(prs.slides) == 2
        # 母片未變
        assert prs.slide_layouts is not None

    def test_no_modify_master(self):
        """確認不破壞母片"""
        from fa_improver.layout.protector import MasterProtector

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        protector = MasterProtector(prs)
        original_xml = prs.slide_masters[0].element.xml

        from fa_improver.domain.evaluation import EvaluationResult

        evaluation = EvaluationResult(total_score=50.0, grade="F")
        add_problem_definition_slide(prs, evaluation)

        protector.verify_unchanged(prs)
        assert prs.slide_masters[0].element.xml == original_xml


class TestAnalysisMethodImprover:
    """分析方法與流程 Improver 測試"""

    def test_add_slide_basic(self):
        """基本新增測試"""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])

        from fa_improver.domain.evaluation import EvaluationResult

        evaluation = EvaluationResult(total_score=50.0, grade="F")
        add_analysis_method_slide(prs, evaluation)

        assert len(prs.slides) == 2

    def test_8d_steps_in_output(self):
        """確認 8D 步驟都包含"""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])

        from fa_improver.domain.evaluation import EvaluationResult

        evaluation = EvaluationResult(total_score=50.0, grade="F")
        add_analysis_method_slide(prs, evaluation)

        # 收集所有文字
        all_text = ""
        for shape in prs.slides[1].shapes:
            if hasattr(shape, "text_frame"):
                all_text += shape.text_frame.text

        # 應包含 D1-D8 關鍵字
        for step in ["D1", "D2", "D3", "D4", "D5", "D6", "D7", "D8"]:
            assert step in all_text, f"應包含 {step}"


class TestEvidenceChecklistImprover:
    """數據與證據支持 Improver 測試"""

    def test_add_slide_basic(self):
        """基本新增測試"""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])

        from fa_improver.domain.evaluation import EvaluationResult

        evaluation = EvaluationResult(total_score=50.0, grade="F")
        add_evidence_checklist_slide(prs, evaluation)

        assert len(prs.slides) == 2

    def test_no_modify_master(self):
        """確認不破壞母片"""
        from fa_improver.layout.protector import MasterProtector

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        protector = MasterProtector(prs)
        original_xml = prs.slide_masters[0].element.xml

        from fa_improver.domain.evaluation import EvaluationResult

        evaluation = EvaluationResult(total_score=50.0, grade="F")
        add_evidence_checklist_slide(prs, evaluation)

        protector.verify_unchanged(prs)
        assert prs.slide_masters[0].element.xml == original_xml


class TestOrchestratorIntegration:
    """Orchestrator 整合測試"""

    def test_plan_triggers_all_six_dimensions(self):
        """測試 6 個維度都能觸發對應動作"""
        from fa_improver.domain.evaluation import Dimension, DimensionScore, EvaluationResult
        from fa_improver.improvers.orchestrator import (
            ImprovementOrchestrator,
            SlideAction,
        )

        # 構造低分評估(觸發所有 6 個維度)
        dimensions = [
            DimensionScore(name=Dimension.BASIC_INFO, score=40, weight=15, comment="缺批號"),
            DimensionScore(name=Dimension.PROBLEM_DEF, score=45, weight=15, comment="未量化"),
            DimensionScore(name=Dimension.METHOD, score=50, weight=20, comment="缺深度分析"),
            DimensionScore(name=Dimension.EVIDENCE, score=45, weight=20, comment="缺對照組"),
            DimensionScore(name=Dimension.ROOT_CAUSE, score=40, weight=20, comment="推測非分析"),
            DimensionScore(name=Dimension.PREVENTION, score=40, weight=10, comment="缺失"),
        ]
        evaluation = EvaluationResult(
            total_score=43.5,
            grade="F",
            dimensions=dimensions,
        )

        from pathlib import Path

        # 用一個 fixture pptx
        pptx_path = Path(__file__).parent.parent.parent / "report" / "MS_Meishan_ADO_445239_260716.pptx"
        if not pptx_path.exists():
            pytest.skip("範例 pptx 不存在")

        orchestrator = ImprovementOrchestrator(evaluation, pptx_path)
        plan = orchestrator.build_plan()

        action_values = {a.value for a in plan.actions}

        # 應包含新增的 3 個動作
        assert "add_problem_definition" in action_values
        assert "add_analysis_method" in action_values
        assert "add_evidence_checklist" in action_values

    def test_high_scores_skip_actions(self):
        """高分時不觸發對應動作"""
        from fa_improver.domain.evaluation import Dimension, DimensionScore, EvaluationResult
        from fa_improver.improvers.orchestrator import (
            ImprovementOrchestrator,
            SlideAction,
        )

        # 構造滿分評估
        dimensions = [
            DimensionScore(name=d, score=95, weight=15 if d in [Dimension.BASIC_INFO, Dimension.PROBLEM_DEF] else 20)
            for d in [
                Dimension.BASIC_INFO,
                Dimension.PROBLEM_DEF,
                Dimension.METHOD,
                Dimension.EVIDENCE,
                Dimension.ROOT_CAUSE,
            ]
        ]
        dimensions.append(DimensionScore(name=Dimension.PREVENTION, score=90, weight=10))
        evaluation = EvaluationResult(
            total_score=92.0,
            grade="A",
            dimensions=dimensions,
        )

        from pathlib import Path

        pptx_path = Path(__file__).parent.parent.parent / "report" / "MS_Meishan_ADO_445239_260716.pptx"
        if not pptx_path.exists():
            pytest.skip("範例 pptx 不存在")

        orchestrator = ImprovementOrchestrator(evaluation, pptx_path)
        plan = orchestrator.build_plan()

        action_values = {a.value for a in plan.actions}

        # 高分時,不應觸發改善動作(只有 Summary 強化)
        assert "add_problem_definition" not in action_values
        assert "add_analysis_method" not in action_values
        assert "add_evidence_checklist" not in action_values
        assert "add_basic_info" not in action_values
        # 但 Summary 強化永遠執行
        assert "enhance_summary" in action_values