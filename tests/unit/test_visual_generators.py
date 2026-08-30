"""視覺元素生成器測試"""

import pytest
from pptx import Presentation
from pptx.util import Inches

from fa_improver.visuals import (
    ELAN_BLUE,
    ELAN_GREEN,
    ChecklistGenerator,
    ComparisonTableGenerator,
    FlowDiagramGenerator,
    ProgressBarGenerator,
    TimelineGenerator,
)


@pytest.fixture
def empty_slide():
    """建立空投影片"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    return slide, prs


class TestChecklistGenerator:
    """Checklist 生成器測試"""

    def test_generate_with_string_list(self, empty_slide):
        """字串列表輸入"""
        slide, _ = empty_slide
        gen = ChecklistGenerator(slide, 0.5, 0.5, 4.0, 2.0)
        gen.generate(["項目1", "項目2", "項目3"])

        # 應新增 2 個 shape per item(checkbox + textbox)
        assert len(slide.shapes) >= 6

    def test_generate_with_dict_list(self, empty_slide):
        """dict 列表輸入(含狀態與顏色)"""
        slide, _ = empty_slide
        gen = ChecklistGenerator(slide, 0.5, 0.5, 4.0, 2.0)
        gen.generate(
            [
                {"text": "已完成", "checked": True, "color": ELAN_GREEN},
                {"text": "未完成", "checked": False, "color": ELAN_BLUE},
            ]
        )
        assert len(slide.shapes) >= 4

    def test_empty_content(self, empty_slide):
        """空內容不應崩潰"""
        slide, _ = empty_slide
        gen = ChecklistGenerator(slide, 0.5, 0.5, 4.0, 2.0)
        gen.generate([])
        # 不應新增任何 shape
        assert len(slide.shapes) == 0


class TestFlowDiagramGenerator:
    """流程圖生成器測試"""

    def test_generate_simple_flow(self, empty_slide):
        """簡單流程"""
        slide, _ = empty_slide
        gen = FlowDiagramGenerator(slide, 0.5, 0.5, 6.0, 1.0)
        gen.generate(["Step 1", "Step 2", "Step 3"])

        # 每個 step 有 box,加上 step-1 個 arrow
        assert len(slide.shapes) >= 3 + 2

    def test_flow_with_status(self, empty_slide):
        """含狀態的流程(用於 5-Why)"""
        slide, _ = empty_slide
        gen = FlowDiagramGenerator(slide, 0.5, 0.5, 6.0, 1.0)
        gen.generate(
            [
                {"name": "Why 1: 通訊失敗", "status": "done"},
                {"name": "Why 2: 訊號異常", "status": "done"},
                {"name": "Why 3: I/O 損傷", "status": "current"},
                {"name": "Why 4: ⚠️ 待驗證", "status": "missing"},
                {"name": "Why 5: ESD?", "status": "pending"},
            ]
        )
        assert len(slide.shapes) >= 5

    def test_single_step(self, empty_slide):
        """單一步驟(無箭頭)"""
        slide, _ = empty_slide
        gen = FlowDiagramGenerator(slide, 0.5, 0.5, 6.0, 1.0)
        gen.generate(["Only Step"])
        # 只有 box,沒有 arrow
        assert len(slide.shapes) == 1


class TestComparisonTableGenerator:
    """對照表生成器測試"""

    def test_generate_table(self, empty_slide):
        """簡單對照表"""
        slide, _ = empty_slide
        gen = ComparisonTableGenerator(slide, 0.5, 0.5, 5.0, 2.0)
        gen.generate(
            [
                ["DVT", "PVT"],
                ["5", "5"],
                ["正常", "異常"],
            ]
        )

        # 應有一個 table shape
        table_count = sum(1 for s in slide.shapes if s.has_table)
        assert table_count == 1

    def test_generate_with_dict(self, empty_slide):
        """用 dict 格式"""
        slide, _ = empty_slide
        gen = ComparisonTableGenerator(slide, 0.5, 0.5, 5.0, 2.0)
        gen.generate(
            {
                "headers": ["參數", "DVT", "PVT"],
                "rows": [
                    ["VH/VOUT", "正常", "異常"],
                    ["ESD", "通過", "失敗"],
                ],
            }
        )
        table_count = sum(1 for s in slide.shapes if s.has_table)
        assert table_count == 1


class TestProgressBarGenerator:
    """進度條生成器測試"""

    def test_generate_basic(self, empty_slide):
        """基本進度條"""
        slide, _ = empty_slide
        gen = ProgressBarGenerator(slide, 0.5, 0.5, 5.0, 2.0)
        gen.generate(
            [
                {"label": "基本資訊", "value": 70, "max_value": 100},
                {"label": "根因分析", "value": 50, "max_value": 100},
                {"label": "改善對策", "value": 70, "max_value": 100},
            ]
        )

        # 每個進度條有 3 個 shape:label, bg, fg + value textbox
        assert len(slide.shapes) >= 12

    def test_color_by_score(self, empty_slide):
        """分數對應顏色"""
        slide, _ = empty_slide
        gen = ProgressBarGenerator(slide, 0.5, 0.5, 5.0, 2.0)
        gen.generate(
            [
                {"label": "高(綠)", "value": 90, "max_value": 100, "color": ELAN_GREEN},
                {"label": "中(藍)", "value": 75, "max_value": 100, "color": ELAN_BLUE},
            ]
        )
        assert len(slide.shapes) >= 6

    def test_empty_content(self, empty_slide):
        """空內容"""
        slide, _ = empty_slide
        gen = ProgressBarGenerator(slide, 0.5, 0.5, 5.0, 2.0)
        gen.generate([])
        assert len(slide.shapes) == 0


class TestTimelineGenerator:
    """時間軸生成器測試"""

    def test_generate_timeline(self, empty_slide):
        """基本時間軸"""
        slide, _ = empty_slide
        gen = TimelineGenerator(slide, 0.5, 0.5, 8.0, 1.5)
        gen.generate(
            [
                {"label": "立即", "timeframe": "本週", "color": ELAN_BLUE},
                {"label": "短期", "timeframe": "1個月", "color": ELAN_GREEN},
                {"label": "中期", "timeframe": "本季", "color": ELAN_BLUE},
                {"label": "長期", "timeframe": "年度", "color": ELAN_GREEN},
            ]
        )

        # 每階段有 node + label + timeframe + 1 個 axis
        assert len(slide.shapes) >= 13

    def test_simplified_string_list(self, empty_slide):
        """簡化字串列表"""
        slide, _ = empty_slide
        gen = TimelineGenerator(slide, 0.5, 0.5, 8.0, 1.5)
        gen.generate(["立即", "短期", "中期", "長期"])
        assert len(slide.shapes) >= 9


class TestMasterProtection:
    """確認視覺元素不會破壞母片"""

    def test_generators_preserve_master(self, sample_pptx):
        """使用視覺元素後母片仍保持"""
        if not sample_pptx.exists():
            pytest.skip("範例 pptx 不存在")

        from fa_improver.layout.protector import MasterProtector

        prs = Presentation(sample_pptx)
        protector = MasterProtector(prs)
        original_master_xml = prs.slide_masters[0].element.xml

        slide = prs.slides[0]
        # 加各種視覺元素
        ChecklistGenerator(slide, 1, 1, 4, 2).generate(["a", "b", "c"])
        FlowDiagramGenerator(slide, 1, 3, 6, 1).generate(["s1", "s2"])
        ProgressBarGenerator(slide, 1, 4.5, 5, 1.5).generate(
            [{"label": "x", "value": 80, "max_value": 100}]
        )

        # 母片保護應通過
        protector.verify_unchanged(prs)
        assert prs.slide_masters[0].element.xml == original_master_xml