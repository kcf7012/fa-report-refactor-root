"""母片保護測試 — 最關鍵的測試"""

import pytest
from pptx import Presentation

from fa_improver.layout.protector import MasterProtector, MasterProtectionError


class TestMasterProtection:
    """母片保護機制測試"""

    def test_snapshot_captures_xml(self, sample_pptx):
        """快照應該擷取母片 XML"""
        if not sample_pptx.exists():
            pytest.skip("範例 pptx 不存在")

        prs = Presentation(sample_pptx)
        protector = MasterProtector(prs)

        assert len(protector.snapshot.masters_xml) == len(prs.slide_masters)
        assert len(protector.snapshot.layouts_xml) == len(prs.slide_layouts)
        assert len(protector.snapshot.layout_names) == len(prs.slide_layouts)

    def test_no_modification_passes(self, sample_pptx):
        """未修改母片應通過驗證"""
        if not sample_pptx.exists():
            pytest.skip("範例 pptx 不存在")

        prs = Presentation(sample_pptx)
        protector = MasterProtector(prs)

        # 立即驗證(未做任何修改)
        protector.verify_unchanged(prs)  # 不應拋出例外

    def test_new_slide_preserves_master(self, sample_pptx):
        """新增投影片不應破壞母片"""
        if not sample_pptx.exists():
            pytest.skip("範例 pptx 不存在")

        prs = Presentation(sample_pptx)
        protector = MasterProtector(prs)

        # 新增一張投影片
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE

        new_slide = prs.slides.add_slide(prs.slide_layouts[1])
        new_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))

        # 母片保護應通過
        protector.verify_unchanged(prs)

    def test_layout_name_validation(self, sample_pptx):
        """確認 layout 名稱檢查"""
        if not sample_pptx.exists():
            pytest.skip("範例 pptx 不存在")

        prs = Presentation(sample_pptx)
        protector = MasterProtector(prs)

        existing_names = [l.name for l in prs.slide_layouts]

        # 存在的 layout 應該可以
        if existing_names:
            protector.assert_can_add_slide(prs, existing_names[0])

        # 不存在的 layout 應該拋出例外
        with pytest.raises(MasterProtectionError):
            protector.assert_can_add_slide(prs, "FakeLayoutName12345")